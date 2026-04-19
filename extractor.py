import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
import os

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        return extract_from_pdf(file_path)
    elif ext in ['.png', '.jpg', '.jpeg', '.webp']:
        return extract_from_image(file_path)
    elif ext == '.docx':
        return extract_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_from_pptx(file_path)
    else:
        return None

def extract_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text.strip()

def extract_from_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text.strip()

def extract_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for para in doc.paragraphs:
        if para.text.strip():
            text += para.text.strip() + "\n"
    return text.strip()

def extract_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
    return text.strip()